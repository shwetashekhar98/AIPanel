"""Tool that parses a PowerPoint (.pptx) file and extracts structured content."""

from typing import Type

from crewai.tools import BaseTool
from pydantic import BaseModel, Field


class PPTXAnalysisInput(BaseModel):
    file_path: str = Field(..., description="Absolute path to the .pptx file")


class PPTXAnalysisTool(BaseTool):
    name: str = "PowerPoint Analyzer"
    description: str = (
        "Parses a .pptx file and extracts all slide content including titles, "
        "body text, speaker notes, image descriptions, table data, and slide count. "
        "Returns a structured text representation of the entire deck."
    )
    args_schema: Type[BaseModel] = PPTXAnalysisInput

    def _run(self, file_path: str) -> str:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        try:
            prs = Presentation(file_path)
            parts: list[str] = []
            parts.append(f"## PowerPoint Analysis")
            parts.append(f"Total slides: {len(prs.slides)}")

            width = prs.slide_width
            height = prs.slide_height
            if width and height:
                parts.append(f"Slide dimensions: {width} x {height} EMUs")

            for idx, slide in enumerate(prs.slides, 1):
                parts.append(f"\n### Slide {idx}")

                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"Speaker Notes: {notes}")

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            if shape == slide.shapes.title:
                                parts.append(f"**Title:** {text}")
                            else:
                                parts.append(f"Text: {text}")

                    if shape.has_table:
                        table = shape.table
                        parts.append("Table:")
                        for row_idx, row in enumerate(table.rows):
                            cells = [cell.text.strip() for cell in row.cells]
                            parts.append(f"  Row {row_idx}: {' | '.join(cells)}")

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        parts.append(f"[Image: {shape.name}, {shape.width}x{shape.height} EMUs]")

                    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        parts.append(f"[Chart: {shape.name}]")

            return "\n".join(parts)

        except Exception as e:
            return f"Error parsing PowerPoint file: {str(e)}"
